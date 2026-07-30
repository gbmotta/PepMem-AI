#!/usr/bin/env python3
"""Gera PPTX editável da apresentação PepMem-AI (texto real + identidade InovAI)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "docs" / "pipeline" / "template_assets"
OUT = ROOT / "docs" / "pipeline" / "PepMem_AI_Modelagem_Dados_Resultados.pptx"

# Paleta InovAI (MODELO DE SLIDE)
PURPLE = RGBColor(0x3C, 0x23, 0xA3)
PURPLE_SOFT = RGBColor(0x85, 0x8D, 0xFD)
LILAC = RGBColor(0xA6, 0x9F, 0xFD)
BLUE = RGBColor(0x42, 0x85, 0xF4)
CYAN = RGBColor(0x00, 0x97, 0xA7)
ORANGE = RGBColor(0xFF, 0xAB, 0x40)
INK = RGBColor(0x21, 0x21, 0x21)
MUTED = RGBColor(0x78, 0x90, 0x9C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xEE, 0xEE, 0xEE)
CARD_BG = RGBColor(0xF7, 0xF6, 0xFC)


def _set_run(run, text: str, size_pt: float, bold: bool = False, color: RGBColor = INK, font: str = "Calibri"):
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return box, tf


def write_lines(tf, lines: list[tuple], clear_first=True):
    """lines: list of (text, size, bold, color, align?, space_after?)."""
    if clear_first:
        tf.clear()
    first = True
    for item in lines:
        text, size, bold, color = item[0], item[1], item[2], item[3]
        align = item[4] if len(item) > 4 else PP_ALIGN.LEFT
        space_after = item[5] if len(item) > 5 else 8
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        run = p.add_run()
        _set_run(run, text, size, bold, color)


def add_bg(slide, path: Path):
    if path.exists():
        slide.shapes.add_picture(str(path), Emu(0), Emu(0), width=Inches(13.333), height=Inches(7.5))


def add_footer(slide, page: int, total: int = 19):
    # white strip
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.95), Inches(13.333), Inches(0.55))
    strip.fill.solid()
    strip.fill.fore_color.rgb = WHITE
    strip.line.fill.background()

    logo_i = ASSETS / "logo_inovai.png"
    logo_u = ASSETS / "logo_ufrn.png"
    if logo_i.exists():
        slide.shapes.add_picture(str(logo_i), Inches(0.35), Inches(7.05), height=Inches(0.35))
    if logo_u.exists():
        slide.shapes.add_picture(str(logo_u), Inches(11.55), Inches(7.05), height=Inches(0.32))

    box, tf = textbox(slide, Inches(5.5), Inches(7.05), Inches(2.3), Inches(0.35))
    write_lines(tf, [("PepMem-AI", 11, True, PURPLE, PP_ALIGN.CENTER, 0)])

    box, tf = textbox(slide, Inches(12.55), Inches(7.08), Inches(0.7), Inches(0.3))
    write_lines(tf, [(f"{page}/{total}", 10, False, MUTED, PP_ALIGN.RIGHT, 0)])


def add_header(slide, title: str):
    box, tf = textbox(slide, Inches(0.55), Inches(0.25), Inches(12), Inches(0.35))
    write_lines(tf, [("INOVAI LAB  ·  UFRN  ·  TITYUS STIGMURUS", 10, True, MUTED, PP_ALIGN.LEFT, 2)])

    box, tf = textbox(slide, Inches(0.55), Inches(0.55), Inches(12), Inches(0.55))
    write_lines(tf, [(title, 28, True, PURPLE, PP_ALIGN.LEFT, 4)])

    # accent bars
    for left, color, w in [
        (0.55, PURPLE, 1.0),
        (1.55, PURPLE_SOFT, 0.35),
        (1.95, BLUE, 0.22),
    ]:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(1.15), Inches(w), Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()


def add_card(slide, left, top, width, height, title: str, body_lines: list[str]):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = LINE

    box, tf = textbox(slide, left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(0.35))
    write_lines(tf, [(title, 14, True, PURPLE, PP_ALIGN.LEFT, 4)])

    box, tf = textbox(slide, left + Inches(0.15), top + Inches(0.45), width - Inches(0.3), height - Inches(0.55))
    lines = [(ln, 13, False, INK, PP_ALIGN.LEFT, 6) for ln in body_lines]
    write_lines(tf, lines)


def add_bullets(slide, left, top, width, height, items: list[str], size=16):
    box, tf = textbox(slide, left, top, width, height)
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        p.level = 0
        run = p.add_run()
        _set_run(run, f"▸  {item}", size, False, INK)


def add_numbered(slide, left, top, width, height, items: list[str], size=16):
    box, tf = textbox(slide, left, top, width, height)
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        run = p.add_run()
        _set_run(run, f"{i + 1}.  {item}", size, False, INK)


def metric_card(slide, left, top, label: str, value: str, note: str, value_color=CYAN):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.5), Inches(1.55))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE

    box, tf = textbox(slide, left + Inches(0.15), top + Inches(0.12), Inches(3.2), Inches(1.3))
    write_lines(
        tf,
        [
            (label, 11, False, MUTED, PP_ALIGN.CENTER, 4),
            (value, 32, True, value_color, PP_ALIGN.CENTER, 4),
            (note, 11, False, MUTED, PP_ALIGN.CENTER, 0),
        ],
    )


def content_slide(prs, title: str, page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, ASSETS / "bg_content_white.png")
    add_header(slide, title)
    add_footer(slide, page)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ---- 1 Título ----
    s = prs.slides.add_slide(blank)
    add_bg(s, ASSETS / "bg_dark_purple.png")
    box, tf = textbox(s, Inches(0.7), Inches(1.2), Inches(10), Inches(5))
    write_lines(
        tf,
        [
            ("INOVAI LAB  ·  LANCE  ·  IMD/UFRN", 12, True, LILAC, PP_ALIGN.LEFT, 18),
            ("PepMem-AI", 44, True, WHITE, PP_ALIGN.LEFT, 12),
            ("Dados, modelagem e resultados", 24, False, WHITE, PP_ALIGN.LEFT, 16),
            ("Priorização in silico de peptídeos escorpiônicos", 16, False, WHITE, PP_ALIGN.LEFT, 4),
            ("antes da validação experimental in vitro", 16, False, WHITE, PP_ALIGN.LEFT, 18),
            ("PoC computacional — baseline RF + multimodal RF/ESM-2", 13, False, LILAC, PP_ALIGN.LEFT, 8),
            ("Métricas: data/processed/models/", 11, False, PURPLE_SOFT, PP_ALIGN.LEFT, 0),
        ],
    )
    logo = ASSETS / "logo_inovai_on_purple.png"
    if logo.exists():
        s.shapes.add_picture(str(logo), Inches(0.7), Inches(6.55), height=Inches(0.55))

    # ---- 2 Agenda ----
    s = content_slide(prs, "Agenda", 2)
    add_numbered(
        s,
        Inches(0.7),
        Inches(1.45),
        Inches(11.5),
        Inches(5),
        [
            "Problema e pergunta do modelo",
            "Dados: fontes, volume e rótulo",
            "Descritores: carga, hidrofobicidade e μH",
            "PMI: fórmula e origem (in silico)",
            "Features: baseline e multimodal (ESM-2)",
            "Modelagem: RF, LOPO e calibração",
            "Resultados (baseline e multimodal)",
            "Uso: dashboard, API e deploy",
            "Limitações e próximos passos",
        ],
        size=16,
    )

    # ---- 3 Problema ----
    s = content_slide(prs, "Problema", 3)
    add_card(
        s,
        Inches(0.55),
        Inches(1.45),
        Inches(12.2),
        Inches(1.5),
        "Contexto",
        [
            "Peptídeos de Tityus stigmurus (Stigmurin e análogos) interagem com membranas diversas",
            "(Gram+/Gram−, fungo, célula mamífera). Testar todos os pares in vitro é caro.",
        ],
    )
    add_card(
        s,
        Inches(0.55),
        Inches(3.2),
        Inches(7.2),
        Inches(2.8),
        "Pergunta operacional",
        [
            "• Dado peptídeo + membrana-alvo,",
            "• qual a chance de alta atividade?",
            "• Rótulo: MIC ≤ 3,4 µM",
        ],
    )
    add_card(
        s,
        Inches(8.0),
        Inches(3.2),
        Inches(4.75),
        Inches(2.8),
        "Papel do PoC",
        [
            "Priorizar candidatos.",
            "Não substitui MIC,",
            "hemólise nem clínica.",
        ],
    )

    # ---- 4 De onde veio ----
    s = content_slide(prs, "De onde veio o desenho", 4)
    add_card(
        s,
        Inches(0.55),
        Inches(1.45),
        Inches(6.0),
        Inches(2.2),
        "Pipeline InovAI (Nível 1)",
        [
            "Baselines: RF, XGBoost, SVM, logreg.",
            "Com poucos dados: classificador +",
            "encoder congelado; PMI + SHAP.",
        ],
    )
    add_card(
        s,
        Inches(0.55),
        Inches(3.9),
        Inches(6.0),
        Inches(2.2),
        "PoC implementado",
        [
            "RF (baseline e multimodal).",
            "ESM-2 t6_8M como embedding.",
            "SHAP TreeExplainer no dashboard.",
        ],
    )
    add_card(
        s,
        Inches(6.9),
        Inches(1.45),
        Inches(5.85),
        Inches(4.65),
        "Proposta CNPq (horizonte)",
        [
            "Deep learning multimodal",
            "(MLP / Transformers).",
            "",
            "Meta de médio/longo prazo;",
            "não é o modelo em produção no Cloud.",
        ],
    )

    # ---- 5 Dados fontes ----
    s = content_slide(prs, "Dados — fontes", 5)
    # flow boxes
    labels = [
        (0.55, "Literatura\nParente / Amorim"),
        (4.55, "Bancada\nmic_bench.csv"),
        (8.55, "OPM / APD\n(bases públicas)"),
    ]
    for x, lab in labels:
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.55), Inches(3.2), Inches(1.1))
        sh.fill.solid()
        sh.fill.fore_color.rgb = WHITE
        sh.line.color.rgb = PURPLE
        box, tf = textbox(s, Inches(x + 0.1), Inches(1.7), Inches(3.0), Inches(0.9))
        write_lines(tf, [(lab.replace("\n", " · "), 13, True, PURPLE, PP_ALIGN.CENTER, 0)])

    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.0), Inches(3.1), Inches(5.2), Inches(1.0))
    sh.fill.solid()
    sh.fill.fore_color.rgb = WHITE
    sh.line.color.rgb = ORANGE
    box, tf = textbox(s, Inches(4.15), Inches(3.25), Inches(4.9), Inches(0.75))
    write_lines(tf, [("Pares peptídeo × alvo + PMI + MIC", 14, True, ORANGE, PP_ALIGN.CENTER, 0)])

    add_bullets(
        s,
        Inches(0.7),
        Inches(4.4),
        Inches(11.5),
        Inches(2.2),
        [
            "Bancada sobrescreve literatura no mesmo par (import_bench_mic.py).",
            "Distribuição atual dos MICs: 78 bancada + 12 literature = 90.",
            "Alvos com MIC: 18 membranas distintas.",
        ],
    )

    # ---- 6 Volume e rótulo ----
    s = content_slide(prs, "Dados — volume e rótulo", 6)
    add_card(
        s,
        Inches(0.55),
        Inches(1.5),
        Inches(5.9),
        Inches(3.2),
        "Conjunto de treino",
        [
            "• 90 pares com MIC",
            "• 10 peptídeos distintos",
            "• Taxa positiva: 44,4%",
        ],
    )
    add_card(
        s,
        Inches(6.8),
        Inches(1.5),
        Inches(5.9),
        Inches(3.2),
        "Rótulo binário",
        [
            "y = 1  se MIC ≤ 3,4 µM",
            "y = 0  caso contrário",
            "",
            "Limiar operacional do projeto",
            "(valores ativos em Parente 2022).",
        ],
    )
    box, tf = textbox(s, Inches(0.7), Inches(5.1), Inches(12), Inches(0.6))
    write_lines(
        tf,
        [("Mediana dos MICs ≈ 4,7 µM — o corte 3,4 µM não é a mediana.", 14, False, MUTED, PP_ALIGN.CENTER, 0)],
    )

    # ---- 7 Descritores do peptídeo ----
    s = content_slide(prs, "Descritores do peptídeo — da sequência", 7)
    add_card(
        s,
        Inches(0.55),
        Inches(1.4),
        Inches(6.0),
        Inches(2.35),
        "Carga líquida q (pH ~7)",
        [
            "Soma dos resíduos ionizáveis:",
            "K, R = +1;  H = +0,1;",
            "D, E = −1;  demais = 0.",
            "Se houver net_charge anotada, ela prevalece.",
        ],
    )
    add_card(
        s,
        Inches(6.8),
        Inches(1.4),
        Inches(5.95),
        Inches(2.35),
        "Hidrofobicidade h",
        [
            "Média Kyte–Doolittle por resíduo:",
            "h = (1/L) · Σ H(AA_i)",
            "Ex.: I=4,5; L=3,8; K=−3,9; R=−4,5.",
        ],
    )
    add_card(
        s,
        Inches(0.55),
        Inches(3.95),
        Inches(6.0),
        Inches(2.3),
        "Momento hidrofóbico μH",
        [
            "Eisenberg (hélice α, ângulo 100°).",
            "Mede anfifilicidade (faces hidro‑",
            "fóbica / hidrofílica na hélice).",
        ],
    )
    add_card(
        s,
        Inches(6.8),
        Inches(3.95),
        Inches(5.95),
        Inches(2.3),
        "Importante",
        [
            "Tudo é in silico a partir da sequência.",
            "Não é medido na bancada.",
            "Código: peptide_utils.py / pmi.py",
        ],
    )

    # ---- 8 Membrana e PMI ----
    s = content_slide(prs, "Membrana e PMI — como se obtém", 8)
    add_card(
        s,
        Inches(0.55),
        Inches(1.4),
        Inches(6.0),
        Inches(4.8),
        "Descritores da membrana",
        [
            "Tabela membrane_targets:",
            "• q_m (carga superficial)",
            "• Fração aniônica, LPS, peptidoglicano",
            "• Colesterol, ergosterol, envelope viral",
            "",
            "Tipologia OPM + regras do projeto",
            "(Gram+/Gram−/fungo/célula normal…).",
            "h_m na inferência: proxy fixo 0,5.",
        ],
    )
    add_card(
        s,
        Inches(6.8),
        Inches(1.4),
        Inches(5.95),
        Inches(2.5),
        "PMI (não é ensaio)",
        [
            "PMI = α·q_p|q_m| + β·h_p h_m",
            "    + γ·μH_p − δ·col_m",
            "α=1; β=0,5; γ=0,3; δ=0,4",
            "(pesos fixos, não aprendidos no RF).",
        ],
    )
    add_card(
        s,
        Inches(6.8),
        Inches(4.1),
        Inches(5.95),
        Inches(2.1),
        "Bancada vs PMI",
        [
            "Bancada / literatura: MIC (rótulo).",
            "PMI: índice calculado do par.",
            "PMI_sel = PMI_alvo − PMI_normal",
            "(só no ranking, não no RF).",
        ],
    )

    # ---- 9 Features baseline + multimodal ----
    s = content_slide(prs, "Features — baseline (11) e multimodal", 9)
    add_card(
        s,
        Inches(0.55),
        Inches(1.4),
        Inches(6.0),
        Inches(2.5),
        "Baseline — 11 clássicas",
        [
            "Peptídeo: q, h, μH",
            "Membrana: carga, aniônica, colesterol,",
            "LPS, peptidoglicano, ergosterol, viral",
            "+ PMI  ·  Deploy Cloud (sem PyTorch).",
        ],
    )
    add_card(
        s,
        Inches(6.8),
        Inches(1.4),
        Inches(5.95),
        Inches(2.5),
        "Multimodal — + ESM-2",
        [
            "facebook/esm2_t6_8M_UR50D",
            "mean-pool ~320d ⇒ ≈331 features",
            "HF Space / local com torch.",
        ],
    )
    add_card(
        s,
        Inches(0.55),
        Inches(4.15),
        Inches(6.0),
        Inches(2.0),
        "Por que t6 8M?",
        [
            "Só precisa da sequência; cabe em CPU;",
            "complementa o PMI sem estrutura 3D.",
        ],
    )
    add_card(
        s,
        Inches(6.8),
        Inches(4.15),
        Inches(5.95),
        Inches(2.0),
        "SHAP (baseline)",
        [
            "Top: q_p, PMI, carga da membrana.",
            "Multimodal: ESM agregado domina;",
            "PMI ≈ 0 no ranking global.",
        ],
    )

    # ---- 10 Pipeline ----
    s = content_slide(prs, "Pipeline de modelagem", 10)
    steps = ["Pares MIC", "Features", "LOO (ref.)", "LOPO (princ.)", "Isotônica", "RF final", "Artefatos"]
    for i, lab in enumerate(steps):
        x = 0.45 + i * 1.8
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.7), Inches(1.65), Inches(0.95))
        sh.fill.solid()
        sh.fill.fore_color.rgb = WHITE
        colors = [PURPLE, PURPLE, ORANGE, CYAN, BLUE, ORANGE, PURPLE]
        sh.line.color.rgb = colors[i]
        box, tf = textbox(s, Inches(x + 0.05), Inches(1.9), Inches(1.55), Inches(0.65))
        write_lines(tf, [(lab, 11, True, colors[i], PP_ALIGN.CENTER, 0)])
    add_bullets(
        s,
        Inches(0.7),
        Inches(3.2),
        Inches(12),
        Inches(3),
        [
            "Pipeline: StandardScaler → RandomForestClassifier",
            "Baseline: 200 árvores; multimodal: 300 árvores, max_depth=6",
            "class_weight='balanced'",
        ],
    )

    # ---- 11 LOO vs LOPO ----
    s = content_slide(prs, "Validação: LOO vs LOPO", 11)
    add_card(
        s,
        Inches(0.55),
        Inches(1.5),
        Inches(5.9),
        Inches(3.5),
        "LOO (por amostra)",
        [
            "Tira 1 par por vez.",
            "O mesmo peptídeo pode ficar no treino (outro alvo).",
            "",
            "AUC mais otimista",
            "(vazamento por homologia).",
        ],
    )
    add_card(
        s,
        Inches(6.8),
        Inches(1.5),
        Inches(5.9),
        Inches(3.5),
        "LOPO (por peptídeo)",
        [
            "Tira todo o peptídeo.",
            "Generalização para sequência nunca vista.",
            "",
            "Métrica principal +",
            "probs OOF para calibração.",
        ],
    )
    box, tf = textbox(s, Inches(0.7), Inches(5.3), Inches(12), Inches(0.6))
    write_lines(
        tf,
        [("Família Stigmurin: análogos ~94% idênticos ⇒ LOPO é essencial.", 14, True, PURPLE, PP_ALIGN.CENTER, 0)],
    )

    # ---- 12 Calibração ----
    s = content_slide(prs, "Calibração isotônica", 12)
    add_bullets(
        s,
        Inches(0.7),
        Inches(1.5),
        Inches(12),
        Inches(3.2),
        [
            "RF ordena bem, mas predict_proba é voto do ensemble, não frequência real.",
            "Com ~90 MICs, probs brutas tendem a polarizar (0/1).",
            "IsotonicRegression nas probs OOF do LOPO → p_calibrada.",
            "Preserva a ordem dos candidatos; alinha % do dashboard ao LOPO.",
        ],
    )
    add_card(
        s,
        Inches(0.55),
        Inches(4.7),
        Inches(12.2),
        Inches(1.5),
        "Inferência",
        ["p_bruta = RF.predict_proba  →  p_calibrada = calibrator.predict([p_bruta])"],
    )

    # ---- 13 Baseline resultados ----
    s = content_slide(prs, "Resultados — Baseline RF", 13)
    metric_card(s, Inches(0.7), Inches(1.5), "LOPO AUC", "0,88", "métrica principal", CYAN)
    metric_card(s, Inches(4.7), Inches(1.5), "LOO AUC", "0,88", "referência", PURPLE)
    metric_card(s, Inches(8.7), Inches(1.5), "Acurácia LOPO", "0,74", "40 positivos / 90", BLUE)
    add_bullets(
        s,
        Inches(0.7),
        Inches(3.5),
        Inches(12),
        Inches(2.8),
        [
            "Modelo: 11 features clássicas + PMI",
            "Deploy: Streamlit Cloud (sem PyTorch)",
            "Calibração: isotônica sobre OOF LOPO",
            "LOO acc. 0,83 · F1+ (LOPO) 0,73 · Precision+ 0,69 · Recall+ 0,78",
        ],
    )

    # ---- 14 Multimodal resultados ----
    s = content_slide(prs, "Resultados — Multimodal RF + ESM-2", 14)
    metric_card(s, Inches(0.7), Inches(1.5), "LOPO AUC", "0,81", "métrica principal", ORANGE)
    metric_card(s, Inches(4.7), Inches(1.5), "LOO AUC", "0,85", "referência", PURPLE)
    metric_card(s, Inches(8.7), Inches(1.5), "Features", "331", "11 + ESM-2 320d", BLUE)
    add_bullets(
        s,
        Inches(0.7),
        Inches(3.5),
        Inches(12),
        Inches(2.8),
        [
            "Embedding captura similaridade de sequência; no LOPO generaliza um pouco pior nesta família pequena",
            "Acurácia LOPO 0,69 · Precision+ 0,60 · Recall+ 0,90 · F1+ 0,72",
            "Deploy: Hugging Face Space / local com requirements-space.txt",
        ],
    )

    # ---- 15 AUC por peptídeo ----
    s = content_slide(prs, "AUC LOPO por peptídeo (baseline)", 15)
    box, tf = textbox(s, Inches(2.5), Inches(1.6), Inches(8), Inches(3.2))
    write_lines(
        tf,
        [
            ("ID      AUC          ID      AUC", 16, True, PURPLE, PP_ALIGN.CENTER, 12),
            ("P11     0,90         P15     0,83", 16, False, INK, PP_ALIGN.CENTER, 8),
            ("P12     0,95         P16     0,95", 16, False, INK, PP_ALIGN.CENTER, 8),
            ("P13     0,88         P17     0,80", 16, False, INK, PP_ALIGN.CENTER, 8),
            ("P14     0,88         P18     0,67", 16, False, INK, PP_ALIGN.CENTER, 8),
        ],
    )
    add_bullets(
        s,
        Inches(0.7),
        Inches(5.0),
        Inches(12),
        Inches(1.5),
        [
            "P05 e P10: AUC por peptídeo indefinida (uma só classe no fold)",
            "Variabilidade entre análogos — amostra ainda limitada",
        ],
        size=14,
    )

    # ---- 16 Explicabilidade ----
    s = content_slide(prs, "Explicabilidade e produto", 16)
    add_card(s, Inches(0.55), Inches(1.45), Inches(5.9), Inches(2.2), "SHAP", [
        "TreeExplainer no RF.",
        "Beeswarm global + local.",
        "Dims ESM agregadas no gráfico.",
    ])
    add_card(s, Inches(0.55), Inches(3.9), Inches(5.9), Inches(2.2), "Dashboard", [
        "Predição única / lote FASTA",
        "Ranking multi-alvo",
        "Relatórios MD/DOCX/PDF",
    ])
    add_card(s, Inches(6.8), Inches(1.45), Inches(5.9), Inches(2.2), "Onde roda", [
        "Cloud: baseline leve",
        "Space HF: multimodal",
        "API FastAPI (opcional)",
    ])
    add_card(s, Inches(6.8), Inches(3.9), Inches(5.9), Inches(2.2), "Docs", [
        "docs/TREINO.md",
        "docs/INTERPRETACAO_RESULTADOS.md",
    ])

    # ---- 17 Limitações ----
    s = content_slide(prs, "Limitações", 17)
    add_bullets(
        s,
        Inches(0.7),
        Inches(1.5),
        Inches(12),
        Inches(5),
        [
            "Poucos peptídeos distintos (~10 com MIC) — LOPO ajuda, mas amostra é pequena",
            "Limiar 3,4 µM é escolha de projeto, não padrão CLSI",
            "Hemólise existe nos pares, mas o RF classifica atividade MIC, não toxicidade direta",
            "Multimodal ainda é RF+ESM, não o encoder multimodal completo da proposta CNPq",
            "Predição prioriza ensaio; não substitui bancada",
        ],
    )

    # ---- 18 Mensagens-chave ----
    s = content_slide(prs, "Mensagens-chave", 18)
    add_numbered(
        s,
        Inches(0.7),
        Inches(1.5),
        Inches(12),
        Inches(5),
        [
            "Dados: 90 MICs (bancada + literatura), rótulo MIC ≤ 3,4 µM",
            "Descritores: q, h, μH da sequência; membrana tipificada; PMI calculado (não ensaiado)",
            "Modelo: Random Forest (Nível 1 InovAI) + PMI; ESM-2 no multimodal",
            "Validação: LOPO é a métrica honesta; isotônica calibra o dashboard",
            "Resultado: baseline LOPO AUC ≈ 0,88; multimodal LOPO AUC ≈ 0,81",
            "Uso: priorizar candidatos Stigmurin/análogos antes do in vitro",
        ],
        size=15,
    )

    # ---- 19 Obrigado ----
    s = prs.slides.add_slide(blank)
    add_bg(s, ASSETS / "bg_thanks.png")
    box, tf = textbox(s, Inches(1.5), Inches(1.8), Inches(10.3), Inches(3.5))
    write_lines(
        tf,
        [
            ("INOVAI LAB  ·  UFRN", 12, True, WHITE, PP_ALIGN.CENTER, 16),
            ("Obrigado", 44, True, WHITE, PP_ALIGN.CENTER, 12),
            ("Perguntas", 24, False, WHITE, PP_ALIGN.CENTER, 22),
            ("Código: github.com/gbmotta/PepMem-AI", 14, False, WHITE, PP_ALIGN.CENTER, 8),
            ("Artefatos: baseline_mic_loo.json · multimodal_mic_loo.json", 12, False, LILAC, PP_ALIGN.CENTER, 0),
        ],
    )
    logos = ASSETS / "logos_inovai_ufrn.png"
    if logos.exists():
        # white plate behind logos
        plate = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.35), Inches(5.85), Inches(4.6), Inches(0.95))
        plate.fill.solid()
        plate.fill.fore_color.rgb = WHITE
        plate.line.fill.background()
        s.shapes.add_picture(str(logos), Inches(4.55), Inches(5.98), height=Inches(0.7))

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Salvo: {OUT} ({OUT.stat().st_size / 1024 / 1024:.2f} MB, {len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
